# -*- coding: utf-8 -*-
"""Presentation de soutenance — SmartTravel Agency Platform.
Plan : Introduction / Cadre general / Analyse des besoins / Conception /
Realisation / Conclusion. Design soigne + transitions animees."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ---------- Charte graphique ----------
BLUE_DARK = RGBColor(0x0B, 0x2A, 0x4A)
BLUE = RGBColor(0x1E, 0x66, 0xE0)
BLUE_LIGHT = RGBColor(0x3B, 0x9E, 0xFF)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x42, 0x4A, 0x55)
SOFT = RGBColor(0x8A, 0x97, 0xA8)
LIGHT_BG = RGBColor(0xEE, 0xF3, 0xFB)
CARD = RGBColor(0xF7, 0xFA, 0xFE)
BORDER = RGBColor(0xD6, 0xE2, 0xF2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


# ---------- Helpers ----------
def transition(slide, kind="fade", direction=None):
    """Ajoute une transition animee a la diapo."""
    extra = f' dir="{direction}"' if direction else ""
    xml = (
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/'
        'presentationml/2006/main" spd="med">'
        f'<p:{kind}{extra}/></p:transition>'
    )
    slide._element.append(etree.fromstring(xml))


def rect(slide, x, y, w, h, color, rounded=False, border=None, bw=1.0):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if border is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = border
        sp.line.width = Pt(bw)
    sp.shadow.inherit = False
    return sp


def text(slide, x, y, w, h, s, size, color, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, italic=False, font="Calibri", spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        p.line_spacing = spacing
    r = p.add_run()
    r.text = s
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = font
    f.color.rgb = color
    return tb


def bullets(slide, x, y, w, h, items, size=18, color=GREY, gap=10,
            marker="►", mcolor=BLUE):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        r1 = p.add_run()
        r1.text = marker + "  "
        r1.font.size = Pt(size - 3)
        r1.font.name = "Calibri"
        r1.font.bold = True
        r1.font.color.rgb = mcolor
        r2 = p.add_run()
        r2.text = it
        r2.font.size = Pt(size)
        r2.font.name = "Calibri"
        r2.font.color.rgb = color
    return tb


def header(slide, title, kicker=None):
    rect(slide, 0, 0, SW, Inches(1.25), BLUE_DARK)
    rect(slide, 0, Inches(1.25), SW, Inches(0.07), BLUE)
    rect(slide, Inches(0.55), Inches(0.42), Inches(0.12), Inches(0.55), CYAN)
    text(slide, Inches(0.85), Inches(0.18), Inches(11.8), Inches(0.9), title,
         27, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    top = Inches(1.65)
    if kicker:
        text(slide, Inches(0.85), Inches(1.45), Inches(11.6), Inches(0.4),
             kicker, 14, BLUE, bold=True, italic=True)
        top = Inches(2.0)
    return top


def footer_note(slide, note):
    rect(slide, Inches(0.85), Inches(6.5), Inches(11.6), Inches(0.55),
         LIGHT_BG, rounded=True)
    text(slide, Inches(1.1), Inches(6.54), Inches(11.2), Inches(0.46),
         "  " + note, 12.5, BLUE, italic=True, anchor=MSO_ANCHOR.MIDDLE)


def pagenum(slide):
    text(slide, Inches(12.5), Inches(7.02), Inches(0.7), Inches(0.4),
         str(len(prs.slides)), 11, SOFT)


def content(title, items, kicker=None, note=None, size=18):
    s = prs.slides.add_slide(blank)
    top = header(s, title, kicker)
    bullets(s, Inches(0.95), top, Inches(11.5), Inches(4.4), items, size=size)
    if note:
        footer_note(s, note)
    pagenum(s)
    transition(s)
    return s


def section(num, title, subtitle=""):
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, SW, SH, BLUE_DARK)
    rect(s, 0, 0, Inches(0.25), SH, BLUE)
    # gros numero translucide
    text(s, Inches(8.3), Inches(0.2), Inches(5.4), Inches(7.2), num, 320,
         RGBColor(0x14, 0x3A, 0x63), bold=True, align=PP_ALIGN.RIGHT,
         anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(1.0), Inches(3.05), Inches(1.3), Inches(0.10), CYAN)
    text(s, Inches(0.95), Inches(2.2), Inches(8), Inches(0.6),
         "PARTIE " + num, 18, CYAN, bold=True)
    text(s, Inches(0.95), Inches(3.35), Inches(8.5), Inches(1.6), title, 40,
         WHITE, bold=True)
    if subtitle:
        text(s, Inches(0.97), Inches(4.9), Inches(8), Inches(0.8), subtitle,
             17, BLUE_LIGHT, italic=True)
    pagenum(s)
    transition(s, "push", direction="l")
    return s


def cards(title, items, kicker=None, note=None):
    """items = list of (titre, [puces], couleur_accent)."""
    s = prs.slides.add_slide(blank)
    header(s, title, kicker)
    n = len(items)
    gap = Inches(0.4)
    total_w = Inches(11.6)
    cw = (total_w - gap * (n - 1)) // n
    x = Inches(0.85)
    y = Inches(2.05)
    ch = Inches(4.2)
    for (ctitle, cpoints, accent) in items:
        rect(s, x, y, cw, ch, CARD, rounded=True, border=BORDER, bw=1.0)
        rect(s, x, y, cw, Inches(0.14), accent, rounded=False)
        text(s, x + Inches(0.25), y + Inches(0.35), cw - Inches(0.5),
             Inches(0.8), ctitle, 18, BLUE_DARK, bold=True)
        bullets(s, x + Inches(0.25), y + Inches(1.15), cw - Inches(0.45),
                ch - Inches(1.3), cpoints, size=14, gap=8, marker="•",
                mcolor=accent)
        x += cw + gap
    if note:
        footer_note(s, note)
    pagenum(s)
    transition(s)
    return s


def stats(title, items, kicker=None):
    """items = list of (chiffre, libelle)."""
    s = prs.slides.add_slide(blank)
    header(s, title, kicker)
    n = len(items)
    gap = Inches(0.4)
    total_w = Inches(11.6)
    cw = (total_w - gap * (n - 1)) // n
    x = Inches(0.85)
    y = Inches(2.6)
    ch = Inches(2.6)
    palette = [BLUE, CYAN, AMBER, BLUE_LIGHT]
    for i, (big, label) in enumerate(items):
        accent = palette[i % len(palette)]
        rect(s, x, y, cw, ch, BLUE_DARK, rounded=True)
        rect(s, x, y + ch - Inches(0.14), cw, Inches(0.14), accent)
        text(s, x, y + Inches(0.45), cw, Inches(1.2), big, 50, WHITE,
             bold=True, align=PP_ALIGN.CENTER)
        text(s, x, y + Inches(1.75), cw, Inches(0.7), label, 14, BLUE_LIGHT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        x += cw + gap
    pagenum(s)
    transition(s)
    return s


# =================================================================
# SLIDE 1 — COUVERTURE
# =================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SW, SH, BLUE_DARK)
rect(s, 0, 0, Inches(0.3), SH, BLUE)
rect(s, 0, Inches(2.55), SW, Inches(0.10), CYAN)
rect(s, 0, Inches(4.35), SW, Inches(0.03), BLUE_LIGHT)
text(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.5),
     "ESSAT — Departement Genie Logiciel   |   Projet de Fin d'Etudes", 14,
     BLUE_LIGHT, bold=True)
text(s, Inches(0.9), Inches(2.72), Inches(11.5), Inches(1.4),
     "SmartTravel Agency Platform", 46, WHITE, bold=True)
text(s, Inches(0.9), Inches(3.82), Inches(11.5), Inches(0.6),
     "Systeme Multi-Agents de Reservation Touristique Intelligente", 19,
     BLUE_LIGHT, italic=True)
text(s, Inches(0.9), Inches(4.85), Inches(11.5), Inches(0.45),
     "Realise par :  Sarra Ounissi  &  Amira ElFahem", 18, WHITE, bold=True)
text(s, Inches(0.9), Inches(5.35), Inches(11.5), Inches(0.45),
     "Encadre par :  M. Hassan Hachicha", 16, WHITE)
text(s, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.45),
     "Annee universitaire 2025 / 2026", 13, SOFT)
transition(s, "fade")

# =================================================================
# SLIDE 2 — PLAN
# =================================================================
s = prs.slides.add_slide(blank)
header(s, "Plan de la presentation")
plan = [
    ("01", "Introduction"),
    ("02", "Cadre general du projet"),
    ("03", "Analyse & specification des besoins"),
    ("04", "Conception"),
    ("05", "Realisation"),
    ("06", "Conclusion & perspectives"),
]
y = Inches(1.9)
for num, lib in plan:
    rect(s, Inches(0.95), y, Inches(11.4), Inches(0.72), CARD, rounded=True,
         border=BORDER)
    rect(s, Inches(0.95), y, Inches(0.9), Inches(0.72), BLUE, rounded=True)
    text(s, Inches(0.95), y, Inches(0.9), Inches(0.72), num, 20, WHITE,
         bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(2.1), y, Inches(10), Inches(0.72), lib, 18, BLUE_DARK,
         bold=True, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.82)
pagenum(s)
transition(s)

# =================================================================
# PARTIE 1 — INTRODUCTION
# =================================================================
section("01", "Introduction", "Contexte, motivation et organisation du rapport")

content("Introduction generale", [
    "Le secteur du tourisme connait une profonde mutation numerique",
    "Les voyageurs veulent comparer, reserver, payer et etre assistes en quelques clics",
    "De nombreuses agences moyennes fonctionnent encore avec Excel, telephone et e-mail",
    "D'ou un ecart croissant entre les attentes des clients et les pratiques des agences",
    "Notre projet : une plateforme web intelligente couvrant tout le parcours client",
], kicker="Pourquoi ce projet ?")

# =================================================================
# PARTIE 2 — CADRE GENERAL DU PROJET
# =================================================================
section("02", "Cadre general du projet",
        "Organisme d'accueil, problematique, solution et objectifs")

content("Organisme d'accueil — DOXSO TECHNOLOGIES", [
    "Entreprise fondee en 2018 a Sfax, Tunisie",
    "Specialisee en solutions intelligentes basees sur l'IA et les systemes multi-agents",
    "Services : plateformes web, integration IA, paiement en ligne, automatisation",
    "Projet realise en teletravail dans le cadre du stage de fin d'etudes",
], kicker="Notre cadre de stage")

content("Problematique", [
    "Disponibilites non fiables → risque de double reservation",
    "Parcours de reservation long, manuel et source d'erreurs",
    "Offre uniforme : aucune personnalisation",
    "Assistance limitee aux horaires d'ouverture",
    "Marketing entierement manuel et irregulier",
], note="Comment centraliser la reservation tout en exploitant l'IA pour fiabiliser, "
        "personnaliser, assister et automatiser ?")

cards("Solution proposee — SmartTravel Agency Platform", [
    ("Cote client", [
        "Catalogue & filtres",
        "Reservation + conflits",
        "Paiement securise",
        "Historique & factures",
    ], BLUE),
    ("Intelligence IA", [
        "Assistant conversationnel",
        "Co-pilote MCP",
        "Recommandations",
    ], CYAN),
    ("Cote agence", [
        "Tableau de bord admin",
        "Gestion catalogue/clients",
        "Marketing automatise",
    ], AMBER),
], kicker="Une plateforme web complete et intelligente")

content("Objectifs du projet", [
    "Digitaliser le parcours de reservation de bout en bout",
    "Fiabiliser la gestion des disponibilites (detection de conflits)",
    "Securiser les paiements et generer les factures",
    "Assister les utilisateurs 24/7 via un assistant conversationnel",
    "Personnaliser l'experience grace aux recommandations",
    "Outiller l'administrateur (tableau de bord + co-pilote IA)",
    "Automatiser les actions marketing",
], size=16)

# =================================================================
# PARTIE 3 — ANALYSE ET SPECIFICATION DES BESOINS
# =================================================================
section("03", "Analyse & specification des besoins",
        "Acteurs, besoins fonctionnels et non fonctionnels, planification")

cards("Identification des acteurs", [
    ("Visiteur anonyme", [
        "Consulter les offres",
        "Rechercher / filtrer",
        "Reserver & payer (tarif standard)",
    ], BLUE_LIGHT),
    ("Client", [
        "Profil & historique",
        "Reservation -10% fidelite",
        "Avis & recommandations",
        "Assistant IA",
    ], BLUE),
    ("Administrateur", [
        "Catalogue & clients",
        "Reservations & paiements",
        "Promotions & marketing",
        "Co-pilote IA",
    ], AMBER),
], kicker="3 acteurs principaux")

content("Besoins fonctionnels", [
    "Consulter, rechercher, filtrer et trier les offres",
    "Creer, modifier et annuler une reservation",
    "Verifier les disponibilites et effectuer un paiement securise",
    "Gerer le profil, l'historique, les avis et les preferences",
    "Administrer le catalogue, les clients, les promotions et les statistiques",
    "Dialoguer avec l'assistant IA et recevoir des recommandations",
], size=17)

content("Besoins non fonctionnels", [
    "Securite : authentification JWT + validation cote client et serveur",
    "Performance : chargement rapide des pages et des resultats de recherche",
    "Fiabilite : fonctionnement stable et bonne gestion des erreurs",
    "Ergonomie : interfaces intuitives et responsives (desktop, tablette, mobile)",
    "Maintenance : code structure et documente, evolutif",
], size=17)

content("Methodologie : Scrum + Jira", [
    "Demarche agile Scrum : iterations courtes et livraison incrementale",
    "Roles : Product Owner (encadrant), Scrum Master, equipe de dev (binome)",
    "Pilotage avec Jira (cle projet PFES) : backlog, sprints, suivi des tickets",
    "Projet sur ~16 semaines, decoupe en 5 sprints thematiques",
], kicker="Organisation du travail",
   note="Inserer la capture du backlog Jira (fig. 1.2)")

# Timeline des 5 sprints
s = prs.slides.add_slide(blank)
header(s, "Planification : les 5 sprints")
sprints = [
    ("S1", "Inscription &\nAuthentification", BLUE),
    ("S2", "Offres &\nCompte client", BLUE_LIGHT),
    ("S3", "Reservation &\nPaiement", CYAN),
    ("S4", "Administration &\nMarketing", AMBER),
    ("S5", "Recommandation\nintelligente", BLUE),
]
rect(s, Inches(1.0), Inches(3.85), Inches(11.3), Inches(0.06), BORDER)
x = Inches(1.0)
cw = Inches(2.1)
gap = Inches(0.22)
for code, lib, col in sprints:
    rect(s, x, Inches(2.2), cw, Inches(1.5), col, rounded=True)
    text(s, x, Inches(2.45), cw, Inches(0.6), code, 26, WHITE, bold=True,
         align=PP_ALIGN.CENTER)
    text(s, x, Inches(3.05), cw, Inches(0.6), "", 10, WHITE)
    rect(s, x + cw / 2 - Inches(0.06), Inches(3.78), Inches(0.12),
         Inches(0.2), col)
    text(s, x - Inches(0.05), Inches(4.15), cw + Inches(0.1), Inches(1.0),
         lib, 13.5, BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
    x += cw + gap
text(s, Inches(1.0), Inches(5.6), Inches(11.3), Inches(0.5),
     "14 user stories planifiees — taux d'achevement 100 %", 16, BLUE,
     bold=True, italic=True, align=PP_ALIGN.CENTER)
pagenum(s)
transition(s)

# =================================================================
# PARTIE 4 — CONCEPTION
# =================================================================
section("04", "Conception",
        "Modelisation UML et architecture logicielle")

content("Diagramme de cas d'utilisation global", [
    "Formalise les fonctionnalites attendues par chaque acteur",
    "Visiteur, Client et Administrateur autour de la plateforme",
    "Relations include : actions reservees aux utilisateurs authentifies",
], kicker="Vue fonctionnelle",
   note="Inserer le diagramme de cas d'utilisation global (fig. 2.1)")

content("Diagramme de classes", [
    "Structure statique des donnees du systeme",
    "Utilisateur (Client, Administrateur), Offre, Service, Reservation",
    "Paiement, Facture, Avis, Preference et leurs relations",
    "Heritage, associations et compositions entre entites",
], kicker="Vue structurelle",
   note="Inserer le diagramme de classes (fig. 2.4)")

cards("Architecture logicielle en couches", [
    ("Presentation", [
        "React + TypeScript",
        "TailwindCSS",
        "Site public / Client / Admin",
    ], BLUE_LIGHT),
    ("Metier", [
        "API REST Express",
        "Auth, reservation, paiement",
        "Catalogue & IA",
    ], BLUE),
    ("Donnees & IA", [
        "MongoDB (Mongoose)",
        "Serveur MCP, Ollama",
        "n8n, Stripe, Shotstack",
    ], AMBER),
], kicker="Architecture client-serveur 3 couches + IA/automatisation",
   note="Inserer l'architecture en couches (fig. 8.9)")

# =================================================================
# PARTIE 5 — REALISATION
# =================================================================
section("05", "Realisation",
        "Environnement, technologies et fonctionnalites developpees")

cards("Environnement & technologies", [
    ("Front-end", [
        "React + TypeScript",
        "TailwindCSS",
    ], BLUE_LIGHT),
    ("Back-end", [
        "Node.js + Express",
        "MongoDB / Mongoose",
        "JWT, Stripe",
    ], BLUE),
    ("IA & Auto", [
        "Serveur MCP",
        "Ollama (LLM local)",
        "n8n, Shotstack, Buffer",
    ], AMBER),
], kicker="Pile technologique MERN")

content("Sprint 1 & 2 — Acces et catalogue", [
    "Inscription / connexion securisees par JWT + protection des routes",
    "Consultation, recherche, filtrage et detail des offres",
    "Espace client : profil, historique de voyages, avis",
    "Integration de l'assistant conversationnel",
], kicker="Sprints 1 et 2",
   note="Inserer interfaces connexion, offres et profil (fig. 3.5 / 4.6 / 4.8)")

content("Sprint 3 — Reservation & paiement securise", [
    "Creation / modification / annulation d'une reservation",
    "Verification des disponibilites en temps reel",
    "Detection automatique des conflits (anti double reservation)",
    "Paiement securise via Stripe + facture PDF + notifications",
], kicker="Sprint 3",
   note="Inserer le formulaire de reservation + paiement Stripe (fig. 5.6 / 5.8)")

content("Sprint 4 — Administration & co-pilote IA", [
    "Gestion du catalogue, des clients et des promotions",
    "Analyse des ventes et des performances",
    "Co-pilote IA : pilotage de l'agence en langage naturel (MCP)",
    "Campagnes automatisees : n8n + Gmail",
], kicker="Sprint 4  ★ differenciateur",
   note="Inserer l'interface Campagnes promo + schema MCP (fig. 6.5 / 8.1)")

content("Sprint 5 — Recommandation & automatisation video", [
    "Configuration des preferences client (budget, destinations, sejour)",
    "Moteur de scoring pondere : preferences + historique",
    "Offres classees par pertinence avec justification transparente",
    "Generation + publication auto de videos (Shotstack → Buffer → FB/Insta)",
], kicker="Sprint 5  ★ differenciateur",
   note="Inserer Recommandations IA + posts reseaux sociaux (fig. 7.5 / 8.7 / 8.8)")

# =================================================================
# PARTIE 6 — CONCLUSION
# =================================================================
section("06", "Conclusion", "Bilan, chiffres-cles et perspectives")

stats("Chiffres-cles du projet", [
    ("5", "Sprints livres"),
    ("14", "User stories"),
    ("100%", "Achevement"),
    ("~40", "Outils MCP"),
], kicker="Un projet mene a terme")

# Bilan + perspectives (2 colonnes)
s = prs.slides.add_slide(blank)
header(s, "Bilan & perspectives")
rect(s, Inches(0.85), Inches(1.85), Inches(5.6), Inches(4.4), CARD,
     rounded=True, border=BORDER)
rect(s, Inches(0.85), Inches(1.85), Inches(5.6), Inches(0.14), BLUE)
text(s, Inches(1.15), Inches(2.05), Inches(5), Inches(0.5), "Bilan", 20,
     BLUE_DARK, bold=True)
bullets(s, Inches(1.15), Inches(2.7), Inches(5.1), Inches(3.3), [
    "Objectifs fixes atteints",
    "5 sprints livres a 100 %",
    "Chaine complete : reservation → IA → marketing",
    "Competences full-stack, IA et gestion de projet agile",
], size=15, gap=9)
rect(s, Inches(6.9), Inches(1.85), Inches(5.55), Inches(4.4), CARD,
     rounded=True, border=BORDER)
rect(s, Inches(6.9), Inches(1.85), Inches(5.55), Inches(0.14), AMBER)
text(s, Inches(7.2), Inches(2.05), Inches(5), Inches(0.5), "Perspectives", 20,
     BLUE_DARK, bold=True)
bullets(s, Inches(7.2), Inches(2.7), Inches(5.0), Inches(3.3), [
    "Application mobile (React Native)",
    "Notifications WhatsApp / SMS + multilingue",
    "Recommandation par apprentissage (filtrage collaboratif)",
    "Deploiement Docker + CI/CD",
], size=15, gap=9, mcolor=AMBER)
pagenum(s)
transition(s)

# =================================================================
# SLIDE FINALE — MERCI
# =================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SW, SH, BLUE_DARK)
rect(s, 0, 0, Inches(0.3), SH, BLUE)
rect(s, Inches(4.55), Inches(3.95), Inches(4.2), Inches(0.10), CYAN)
text(s, Inches(0.8), Inches(2.7), Inches(11.7), Inches(1.0),
     "Merci de votre attention", 42, WHITE, bold=True, align=PP_ALIGN.CENTER)
text(s, Inches(0.8), Inches(4.25), Inches(11.7), Inches(0.6),
     "Questions & discussion", 22, BLUE_LIGHT, italic=True,
     align=PP_ALIGN.CENTER)
transition(s, "fade")

import sys
out = "SmartTravel_Presentation_v2.pptx"
try:
    prs.save(out)
except PermissionError:
    out = "SmartTravel_Presentation_v3.pptx"
    prs.save(out)
print("OK ->", out, "| slides:", len(prs.slides))
